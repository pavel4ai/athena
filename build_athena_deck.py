"""Build a 3-slide Athena deck (16:9) — readability-first revision.

Larger type throughout; slide 2 uses 2 wide columns x 3 tall rows so titles can
breathe at ~17pt and body at ~13pt without run-on. Theme: Athena Midnight
(navy dominant, gold accent, ice-blue support).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

NAVY   = RGBColor(0x0E, 0x14, 0x33)
NAVY2  = RGBColor(0x16, 0x1F, 0x4D)
INDIGO = RGBColor(0x2B, 0x35, 0x80)
GOLD   = RGBColor(0xE8, 0xB5, 0x4B)
GOLD_D = RGBColor(0xC9, 0x97, 0x2E)
ICE    = RGBColor(0xCA, 0xDC, 0xFC)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CREAM  = RGBColor(0xF6, 0xF7, 0xFB)
INKDK  = RGBColor(0x14, 0x1A, 0x33)
BODYDK = RGBColor(0x3A, 0x42, 0x60)
GREY   = RGBColor(0x9A, 0xA3, 0xC0)
GREEN  = RGBColor(0x3D, 0xDC, 0x84)

HEAD = "Georgia"
BODY = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def bg(slide, color):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color


def rect(slide, x, y, w, h, fill, line=None, rounded=False, line_w=None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(line_w or 1)
    shp.shadow.inherit = False
    return shp


def circle(slide, x, y, d, fill, line=None, line_w=2):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, line_spacing=1.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line_spacing
        for (t, sz, col, bold, ital, fnt) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = bold; r.font.italic = ital; r.font.name = fnt
    return tb


def R(t, sz, col, bold=False, ital=False, fnt=BODY):
    return (t, sz, col, bold, ital, fnt)


# ============================================================ SLIDE 1 — TITLE
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, 0, 0, 0.30, 7.5, GOLD)
rect(s, 8.7, 0, 4.633, 7.5, NAVY2)
circle(s, 10.1, 0.85, 2.05, None, GOLD, 2.5)
circle(s, 10.62, 1.37, 1.0, INDIGO, GOLD, 2)
for cx, cy in [(9.75,0.95),(11.95,1.05),(9.85,3.0),(12.0,2.85),(10.95,3.45)]:
    circle(s, cx, cy, 0.36, GOLD, NAVY, 1.5)
txt(s, 0.95, 1.05, 7.4, 0.6, [[R("ATHENA", 20, GOLD, True, False, BODY)]])
txt(s, 0.9, 1.55, 7.7, 2.2, [
    [R("Agentic Investment", 50, WHITE, True, False, HEAD)],
    [R("Intelligence", 50, WHITE, True, False, HEAD)]], line_spacing=1.0)
txt(s, 0.95, 4.05, 7.5, 1.1, [
    [R("Macro analysis, quantitative reasoning, fundamental", 19, ICE, False, False, BODY)],
    [R("research, and disciplined risk control —", 19, ICE, False, False, BODY)],
    [R("with human-approved execution.", 19, ICE, True, False, BODY)]], line_spacing=1.12)
stats = [("8", "specialist\nAI agents"), ("5", "analytical\nlenses"),
         ("100%", "human-approved\ntrades"), ("Live", "market\ndata")]
x = 0.95
for big, lab in stats:
    txt(s, x, 5.7, 1.95, 1.3, [
        [R(big, 46, GOLD, True, False, HEAD)]], line_spacing=1.0)
    txt(s, x, 6.55, 1.95, 0.7, [[R(l, 13, ICE, False, False, BODY)] for l in lab.split("\n")], line_spacing=1.0, space_after=0)
    x += 2.0

# ====================================================== SLIDE 2 — CAPABILITIES
s = prs.slides.add_slide(blank)
bg(s, CREAM)
rect(s, 0, 0, 13.333, 1.4, NAVY)
rect(s, 0, 1.4, 13.333, 0.07, GOLD)
txt(s, 0.7, 0.30, 9.0, 0.5, [[R("ATHENA", 15, GOLD, True, False, BODY)]])
txt(s, 0.7, 0.62, 11.5, 0.7, [[R("Capabilities & Advantages", 34, WHITE, True, False, HEAD)]])

cards = [
    ("Multi-Agent Research", "Eight specialist agents — Scout, Oracle, Analyst, Allocator, Sentinel, Schwab, Archivist — coordinate through an auditable shared workspace.", GOLD),
    ("Information Alpha", "News and X / Twitter velocity surfaced before institutional pricing — the system's highest-conviction edge.", INDIGO),
    ("Five-Lens Framework", "Every thesis is tested through the Graham, Fabozzi, Hull/McMillan, Chan, and Macro lenses, then scored 0–100.", GOLD),
    ("Disciplined Risk Control", "Independent risk veto, per-position / sector / drawdown limits, and a mandatory exit thesis for every position.", INDIGO),
    ("Human-Approved Execution", "Nothing trades without explicit approval. Athena proposes, you approve, and only then does it execute.", GOLD),
    ("Continuous Learning", "A full audit trail with luck-versus-process attribution and post-mortems after every rebalance.", INDIGO),
]
# 2 columns x 3 rows — wide, tall cards
cw, ch = 5.95, 1.62
gx, gy = 0.45, 0.30
x0, y0 = 0.6, 1.75
for i, (title, desc, accent) in enumerate(cards):
    col = i % 2; row = i // 2
    cx = x0 + col * (cw + gx); cy = y0 + row * (ch + gy)
    rect(s, cx, cy, cw, ch, WHITE, line=RGBColor(0xDD,0xE1,0xEE), rounded=True, line_w=1)
    rect(s, cx, cy, 0.12, ch, accent)
    circle(s, cx+0.34, cy+0.42, 0.62, accent)
    txt(s, cx+0.34, cy+0.42, 0.62, 0.62, [[R(str(i+1), 24, WHITE, True, False, HEAD)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, cx+1.20, cy+0.22, cw-1.5, 0.5, [[R(title, 18, INKDK, True, False, BODY)]])
    txt(s, cx+1.20, cy+0.66, cw-1.5, ch-0.8, [[R(desc, 13, BODYDK, False, False, BODY)]], line_spacing=1.08)

# ====================================================== SLIDE 3 — ARCHITECTURE
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, 0, 0, 0.30, 7.5, GOLD)
txt(s, 0.95, 0.42, 8.0, 0.5, [[R("ATHENA", 15, GOLD, True, False, BODY)]])
txt(s, 0.9, 0.78, 12.0, 0.7, [[R("Architecture & Data Sources", 34, WHITE, True, False, HEAD)]])

# DATA SOURCES (left)
txt(s, 0.9, 1.85, 3.3, 0.4, [[R("DATA SOURCES", 15, GOLD, True, False, BODY)]])
srcs = ["Schwab live market data", "Bloomberg · Reuters · WSJ · CNBC",
        "X / Twitter velocity", "Fed · Treasury · SEC filings",
        "Earnings calls & calendars", "Options / volatility regime",
        "ETF flows · commodities"]
yy = 2.45
for sc in srcs:
    circle(s, 0.95, yy+0.06, 0.18, GOLD)
    txt(s, 1.32, yy-0.02, 3.1, 0.5, [[R(sc, 13, ICE, False, False, BODY)]], anchor=MSO_ANCHOR.MIDDLE)
    yy += 0.62

# AGENT PIPELINE (center)
def node(x, y, w, h, label, sub, fill, lab_col=WHITE, sub_col=ICE):
    rect(s, x, y, w, h, fill, line=GOLD_D, rounded=True, line_w=1.25)
    txt(s, x+0.18, y+0.10, w-0.36, h-0.2, [
        [R(label, 15, lab_col, True, False, BODY)],
        [R(sub, 11, sub_col, False, False, BODY)]], line_spacing=1.0, anchor=MSO_ANCHOR.MIDDLE)

px = 4.85; pw = 3.25; nh = 0.62
chain = [
    ("Athena Scout", "intel · X/news velocity", INDIGO),
    ("Athena Oracle", "theses · regime · 2nd-order", INDIGO),
    ("Athena Analyst", "fundamental + quant", INDIGO),
    ("Athena Allocator", "construction · sizing", NAVY2),
    ("Athena Sentinel", "risk · veto · limits", NAVY2),
]
yy = 1.80; step = 0.82
for idx, (lab, sub, fill) in enumerate(chain):
    node(px, yy, pw, nh, lab, sub, fill)
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(px+pw/2-0.07), Inches(yy+nh+0.01),
                           Inches(0.14), Inches(step-nh-0.02))
    a.fill.solid(); a.fill.fore_color.rgb = GOLD; a.line.fill.background(); a.shadow.inherit=False
    yy += step
node(px, yy, pw, nh+0.06, "★  HUMAN APPROVAL", "Telegram / CLI — required", GOLD,
     lab_col=NAVY, sub_col=NAVY)

# EXECUTION & MEMORY (right)
rx = 8.55; rw = 4.0
txt(s, rx, 1.85, rw, 0.4, [[R("EXECUTION & MEMORY", 15, GOLD, True, False, BODY)]])
rcards = [
    ("Athena Schwab", "Order preview, reconcile, then execute (mock or live).", GREEN),
    ("Athena Archivist", "Decision journal · performance vs SPY / QQQ / BRK.", ICE),
    ("Athena Orchestrator", "Cron cadence · cohort isolation · Telegram briefings.", GOLD),
]
yy = 2.45; rh = 1.30
for lab, desc, accent in rcards:
    rect(s, rx, yy, rw, rh, NAVY2, line=RGBColor(0x33,0x3E,0x70), rounded=True, line_w=1)
    rect(s, rx, yy, 0.11, rh, accent)
    txt(s, rx+0.32, yy+0.18, rw-0.5, 0.4, [[R(lab, 15, WHITE, True, False, BODY)]])
    txt(s, rx+0.32, yy+0.60, rw-0.5, 0.6, [[R(desc, 12, ICE, False, False, BODY)]], line_spacing=1.05)
    yy += rh + 0.18

# (footer removed — kept the bottom clear so nothing collides with the approval box)

import os
out = os.path.expanduser("~/athena_deck.pptx")
prs.save(out)
print("saved:", out)
